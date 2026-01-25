from pptx import Presentation
import os

def read_pptx(file_path):
    if not os.path.exists(file_path):
        print(f"File not found: {file_path}")
        return

    try:
        prs = Presentation(file_path)
        print(f"--- Content of {os.path.basename(file_path)} ---")
        for i, slide in enumerate(prs.slides):
            print(f"\n--- Slide {i+1} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    print(shape.text)
    except Exception as e:
        print(f"Error reading {file_path}: {str(e)}")

if __name__ == "__main__":
    files = [
        r"c:\Users\Tommaso\Documents\HOWEST\TeamProject\MTS3-MCTE-Team-Project-Energy-G1\docs\brief\Introduction week  - Team project.pptx",
        r"c:\Users\Tommaso\Documents\HOWEST\TeamProject\MTS3-MCTE-Team-Project-Energy-G1\docs\brief\Energy.pptx"
    ]
    
    for f in files:
        read_pptx(f)
